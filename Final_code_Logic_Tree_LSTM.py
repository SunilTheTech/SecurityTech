# ---------------------------
# Incognoir: End-to-end Colab Script
# Logic Tree + LSTM Hybrid (Phase1→4)
# Paste into Colab and run (may take GPUs for training)
# ---------------------------

# Install a compatible version of scikit-learn first
!pip install scikit-learn==1.3.0 -q
# Install other packages, including imbalanced-learn
!pip install imbalanced-learn==0.11.0 joblib python-pptx reportlab -q

# Standard imports
import os, json, math, time, joblib
from pathlib import Path
import numpy as np
import pandas as pd
import matplotlib.pyplot as plt
import seaborn as sns
from collections import defaultdict, deque
from sklearn.model_selection import train_test_split, StratifiedKFold
from sklearn.preprocessing import MinMaxScaler, StandardScaler
from sklearn.metrics import (
    accuracy_score, precision_score, recall_score, f1_score,
    roc_auc_score, average_precision_score, confusion_matrix,
    precision_recall_curve, roc_curve
)
from imblearn.over_sampling import SMOTE
import tensorflow as tf
from tensorflow.keras.models import Sequential, load_model
from tensorflow.keras.layers import LSTM, Dense, Dropout, Bidirectional
from tensorflow.keras.callbacks import EarlyStopping, ReduceLROnPlateau, ModelCheckpoint

# Report libs
from pptx import Presentation
from pptx.util import Inches
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Image as RLImage, Table, TableStyle
from reportlab.lib.pagesizes import A4
from reportlab.lib.styles import getSampleStyleSheet
from reportlab.lib import colors

# ---------------------------
# Config & paths
# ---------------------------
RND = 42
np.random.seed(RND)
tf.random.set_seed(RND)

DATA_PATH = "synthetic_golden_dataset_50k.csv"   # change if needed
MODEL_DIR = Path("model_artifacts")
MODEL_DIR.mkdir(exist_ok=True)
SCALER_PATH = MODEL_DIR / "scaler.pkl"
LSTM_MODEL_PATH = MODEL_DIR / "lstm_model.h5"
LOGIC_TREE_PATH = MODEL_DIR / "logic_tree.pkl"
REPORT_PPTX = "LogicTree_LSTM_Report.pptx"
REPORT_PDF = "LogicTree_LSTM_Report.pdf"

# ---------------------------
# 0) Utility: quick dataset check / placeholder
# ---------------------------
if not Path(DATA_PATH).exists():
    print(f"⚠️ {DATA_PATH} not found. Creating small synthetic placeholder dataset for demo.")
    # minimal synthetic dataset to exercise pipeline (replace with real dataset)
    N = 20000
    # columns: session_id, label (0 human,1 bot), and some numeric features similar to your dataset
    df = pd.DataFrame({
        "session_id": [f"sess_{i}" for i in range(N)],
        "cursor_speed_mean": np.random.rand(N),
        "cursor_speed_std": np.random.rand(N),
        "scroll_speed_mean": np.random.rand(N),
        "typing_speed_mean": np.random.rand(N),
        "dwell_time_mean": np.random.rand(N),
        "hesitation_time_mean": np.random.rand(N),
        "click_accuracy": np.random.rand(N),
        "movement_pattern_score": np.random.rand(N),
        "cart_modification_count": np.random.poisson(0.2, N),
        "payment_attempts": np.random.poisson(0.1, N),
        "review_submit_count": np.random.poisson(0.05, N),
        "review_similarity_score": np.random.rand(N),
        "sensitive_field_access": np.random.binomial(1,0.02,N),
    })
    # create imbalanced labels (5% bots)
    df["label"] = np.random.binomial(1, 0.05, size=N)
    df.to_csv(DATA_PATH, index=False)
    print("Placeholder dataset created at", DATA_PATH)

# ---------------------------
# 1) Load dataset & diagnostics
# ---------------------------
print("Loading dataset...")
df = pd.read_csv(DATA_PATH)
print("Original shape:", df.shape)
print(df.head(2))

# Quick diagnostics
print("\n--- Basic diagnostics ---")
print(df.info())
print("\nLabel distribution:\n", df['label'].value_counts(normalize=False))

# ---------------------------
# 2) Feature engineering
# ---------------------------
# Identify numerical columns for model (drop identifiers & obvious leakage if present)
drop_cols = []
for col in ["session_id", "event_id", "user_id_hash", "device_id_hash", "ip_hash", "honeypot_hit", "canary_token_matched", "combined_score", "ml_score"]:
    if col in df.columns: drop_cols.append(col)
X_all = df.drop(columns=drop_cols + ["label"], errors='ignore')
y_all = df["label"].astype(int)

# If any high-cardinality string columns exist, drop for now
for c in X_all.select_dtypes(include='object').columns:
    X_all = X_all.drop(columns=[c])

print("\nSelected features for modeling:", X_all.shape[1])

# Add derived features where possible (example patterns; adapt to your columns)
def add_derived_features(X):
    X = X.copy()
    # if cursor features present make ratios
    if 'cursor_speed_mean' in X.columns and 'scroll_speed_mean' in X.columns:
        X['cursor_scroll_ratio'] = X['cursor_speed_mean'] / (X['scroll_speed_mean'] + 1e-6)
    # movement intensity
    if 'movement_pattern_score' in X.columns:
        X['movement_intensity'] = X['movement_pattern_score'] * X.get('click_accuracy', 1)
    # idle vs activity proxy
    if 'dwell_time_mean' in X.columns and 'typing_speed_mean' in X.columns:
        X['idle_to_typing'] = X['dwell_time_mean'] / (X['typing_speed_mean'] + 1e-6)
    return X

X_all = add_derived_features(X_all)
print("After derived features:", X_all.shape)

# Simple imputation for missing numeric columns
X_all = X_all.fillna(X_all.median())

# ---------------------------
# 3) Scaling & train/test split
# ---------------------------
scaler = MinMaxScaler()
X_scaled = scaler.fit_transform(X_all)
joblib.dump(scaler, SCALER_PATH)
print("Scaler saved to", SCALER_PATH)

# stratified split
X_train, X_test, y_train, y_test = train_test_split(
    X_scaled, y_all, test_size=0.25, random_state=RND, stratify=y_all
)
print("Train/Test shapes:", X_train.shape, X_test.shape, y_train.sum(), y_test.sum())

# ---------------------------
# 4) Class balancing (SMOTE) on training set (only) if imbalanced
# ---------------------------
pos_frac = y_train.mean()
print("Train positive fraction:", pos_frac)
if pos_frac < 0.15:
    print("Applying SMOTE to balance classes on training set...")
    # Removed n_jobs from SMOTE as it's no longer supported in imbalanced-learn 0.11.0+
    sm = SMOTE(random_state=RND)
    X_train_res, y_train_res = sm.fit_resample(X_train, y_train)
    print("After SMOTE:", X_train_res.shape, np.bincount(y_train_res))
else:
    X_train_res, y_train_res = X_train, y_train

# reshape input for LSTM: we will treat features as single time-step sequences (if you have time-series per session, adapt)
# To use LSTM properly, you ideally need sequences (timesteps x features). If only aggregated features exist,
# we use a simple reshape to (samples, timesteps=1, features). For sequence raw data, replace this with actual sequences.
TIMESTEPS = 1
n_features = X_train_res.shape[1]
X_train_lstm = X_train_res.reshape((X_train_res.shape[0], TIMESTEPS, n_features))
X_test_lstm = X_test.reshape((X_test.shape[0], TIMESTEPS, n_features))

print("LSTM input shapes:", X_train_lstm.shape, X_test_lstm.shape)

# ---------------------------
# 5) Build & train enhanced LSTM (Bidirectional + dropout)
# ---------------------------
def build_lstm_model(timesteps, n_features):
    model = Sequential([
        Bidirectional(LSTM(128, return_sequences=False), input_shape=(timesteps, n_features)),
        Dropout(0.35),
        Dense(64, activation='relu'),
        Dropout(0.2),
        Dense(1, activation='sigmoid')
    ])
    model.compile(optimizer='adam', loss='binary_crossentropy', metrics=['accuracy'])
    return model

model = build_lstm_model(TIMESTEPS, n_features)
model.summary()

# Callbacks
callbacks = [
    EarlyStopping(monitor='val_loss', patience=8, restore_best_weights=True, verbose=1),
    ReduceLROnPlateau(monitor='val_loss', factor=0.5, patience=4, verbose=1),
    ModelCheckpoint(str(LSTM_MODEL_PATH), monitor='val_loss', save_best_only=True, verbose=1)
]

EPOCHS = 50
BATCH = 256

history = model.fit(
    X_train_lstm, y_train_res,
    validation_split=0.15,
    epochs=EPOCHS,
    batch_size=BATCH,
    callbacks=callbacks,
    verbose=2
)

# Load best model
if Path(LSTM_MODEL_PATH).exists():
    model = load_model(LSTM_MODEL_PATH)
    print("Loaded saved best model.")

# Save final model
model.save(LSTM_MODEL_PATH)
print("Saved LSTM model to", LSTM_MODEL_PATH)

# Quick training curves
plt.figure(figsize=(8,4))
plt.plot(history.history.get('loss', []), label='train_loss')
plt.plot(history.history.get('val_loss', []), label='val_loss')
plt.legend(); plt.title("Training Loss"); plt.show()

# ---------------------------
# 6) LSTM predictions on test
# ---------------------------
y_proba_lstm = model.predict(X_test_lstm).ravel()
y_pred_lstm = (y_proba_lstm >= 0.5).astype(int)

print("LSTM test metrics (threshold 0.5):")
print("Accuracy:", accuracy_score(y_test, y_pred_lstm))
print("Precision:", precision_score(y_test, y_pred_lstm, zero_division=0))
print("Recall:", recall_score(y_test, y_pred_lstm, zero_division=0))
print("ROC-AUC:", roc_auc_score(y_test, y_proba_lstm))

# ---------------------------
# 7) LogicTree (trie) implementation — fast, in-memory
# ---------------------------
from dataclasses import dataclass, field

@dataclass
class LTNode:
    token: str
    total: int = 0
    malicious: int = 0
    children: dict = field(default_factory=dict)

class LogicTree:
    def __init__(self): # Fixed this line
        self.root = LTNode("<ROOT>")

    def add_path(self, tokens, label):
        cur = self.root
        cur.total += 1
        if label == 1: cur.malicious += 1
        for t in tokens:
            if t not in cur.children:
                cur.children[t] = LTNode(t)
            cur = cur.children[t]
            cur.total += 1
            if label == 1: cur.malicious += 1

    def score_path(self, tokens, alpha=1.0):
        cur = self.root
        risks = []
        for t in tokens:
            if t in cur.children:
                cur = cur.children[t]
                risks.append((cur.malicious + alpha) / (cur.total + 2*alpha))
            else:
                # unseen token -> neutral small risk
                risks.append(alpha/(1+2*alpha))
                break
        return float(np.mean(risks)) if risks else 0.0

    def save(self, path=LOGIC_TREE_PATH):
        joblib.dump(self, path)

    @staticmethod
    def load(path=LOGIC_TREE_PATH):
        if Path(path).exists():
            return joblib.load(path)
        return LogicTree()

logic_tree = LogicTree.load()

# To bootstrap LogicTree from labeled training set we need tokenized paths.
# If you don't have sequence tokens, we can approximate tokens from feature buckets (coarse).
# Example: create tokens from quantized top-k features to simulate path.
def create_synthetic_tokens_from_features(X_raw, n_tokens=6):
    # X_raw is unscaled original features (pandas) or scaled array
    # Here we will quantize top features into tokens to build a path-like sequence
    # This is a fallback if you don't have raw micro-action tokens
    arr = X_raw if isinstance(X_raw, np.ndarray) else X_raw.values
    tokens_list = []
    q = np.percentile(arr, [20,40,60,80], axis=0)
    for row in arr:
        toks = []
        for i, v in enumerate(row[:min(6, arr.shape[1])]):
            # token name = feature_index + bucket
            bucket = np.sum(v > q[:,i]) if q.shape[0] == arr.shape[0] else int(v * 4)
            toks.append(f"f{i}_b{bucket}")
        tokens_list.append(toks)
    return tokens_list

# bootstrap from train set (use original X_train_res features before reshape)
try:
    X_train_for_tokens = X_train_res if isinstance(X_train_res, np.ndarray) else X_train_res.values
    token_paths = create_synthetic_tokens_from_features(X_train_for_tokens)
    for seq, lbl in zip(token_paths, y_train_res):
        logic_tree.add_path(seq, int(lbl))
    logic_tree.save()
    print("Bootstrapped LogicTree from training features and saved.")
except Exception as e:
    print("LogicTree bootstrap skipped:", e)

# ---------------------------
# 8) Hybrid scoring (LSTM + LogicTree)
# ---------------------------
# We need a function that given an example row returns a token path, LSTM prob and combined risk.
# We'll reuse create_synthetic_tokens_from_features for tokens on test set.
X_test_for_tokens = X_test if isinstance(X_test, np.ndarray) else X_test.values
token_paths_test = create_synthetic_tokens_from_features(X_test_for_tokens)

# Combine with weights and compute metrics across a range of thresholds to tune.
w_lstm = 0.6
w_tree = 0.4

combined_scores = []
for proba, path in zip(y_proba_lstm, token_paths_test):
    tree_score = logic_tree.score_path(path)
    combined_score = w_lstm * proba + w_tree * tree_score
    combined_scores.append(combined_score)
combined_scores = np.array(combined_scores)

# Tune threshold by maximizing F1 or meeting FPR <= 0.03 constraint
prec, rec, thresh = precision_recall_curve(y_test, combined_scores)
f1s = 2 * (prec * rec) / (prec + rec + 1e-12)
best_idx = np.nanargmax(f1s)
best_thresh = thresh[best_idx] if len(thresh)>0 else 0.5
print("Best F1 threshold (hybrid):", best_thresh)

# Also find threshold that gives FPR <= 0.03 with highest recall
fpr_vals, tpr_vals, roc_th = roc_curve(y_test, combined_scores)
candidate_thresh = 0.5
# Find the highest threshold that results in FPR <= 0.03
candidate_thresh = None
for t, fprv in zip(reversed(roc_th), reversed(fpr_vals)):
    if fprv <= 0.03:
        candidate_thresh = t
        break
if candidate_thresh is None and len(roc_th) > 0: # If no threshold meets FPR <= 0.03, take the lowest threshold
    candidate_thresh = roc_th[0]


print("Candidate threshold for FPR<=3% (if exists):", candidate_thresh)

# Choose final threshold (tweak as needed)
# Use the candidate_thresh if it exists and is less than or equal to the best_thresh,
# otherwise use the best_thresh. This prioritizes meeting the FPR constraint if possible.
THRESHOLD = candidate_thresh if candidate_thresh is not None and candidate_thresh <= best_thresh else best_thresh

print("Selected hybrid threshold:", THRESHOLD)

y_pred_hybrid = (combined_scores >= THRESHOLD).astype(int)

# ---------------------------
# 9) Metrics & Plots
# ---------------------------
def print_metrics(y_true, y_pred, y_score):
    acc = accuracy_score(y_true, y_pred)
    prec = precision_score(y_true, y_pred, zero_division=0)
    rec = recall_score(y_true, y_pred, zero_division=0)
    f1 = f1_score(y_true, y_pred, zero_division=0)
    roc = roc_auc_score(y_true, y_score) if len(np.unique(y_true))>1 else float('nan')
    pr_auc = average_precision_score(y_true, y_score) if len(np.unique(y_true))>1 else float('nan')
    cm_vals = confusion_matrix(y_true, y_pred).astype(int)
    tn, fp, fn, tp = cm_vals.ravel()
    FPR = fp / (fp + tn + 1e-12)
    print("Accuracy:", round(acc,4), "Precision:", round(prec,4), "Recall:", round(rec,4), "F1:", round(f1,4))
    print("ROC-AUC:", round(roc,4), "PR-AUC:", round(pr_auc,4), "FPR:", round(FPR,4))
    return dict(accuracy=acc, precision=prec, recall=rec, f1=f1, roc=roc, pr_auc=pr_auc, fpr=FPR)

metrics = print_metrics(y_test, y_pred_hybrid, combined_scores)

# Visualizations
plt.figure(figsize=(6,5))
cm = confusion_matrix(y_test, y_pred_hybrid)
sns.heatmap(cm, annot=True, fmt='d', cmap='Blues')
plt.title("Confusion Matrix (Hybrid)")
plt.xlabel("Predicted"); plt.ylabel("True")
plt.show()

# ROC
fpr, tpr, _ = roc_curve(y_test, combined_scores)
plt.figure(figsize=(6,5))
plt.plot(fpr, tpr, label=f"ROC AUC = {metrics['roc']:.3f}")
plt.plot([0,1],[0,1], linestyle='--', color='gray')
plt.xlabel("FPR"); plt.ylabel("TPR"); plt.title("ROC Curve (Hybrid)")
plt.legend(); plt.show()

# PR
prec, rec, _ = precision_recall_curve(y_test, combined_scores)
plt.figure(figsize=(6,5))
plt.plot(rec, prec)
plt.xlabel("Recall"); plt.ylabel("Precision"); plt.title("Precision-Recall (Hybrid)")
plt.show()

# KPI table
df_summary = pd.DataFrame([
    {"Metric":"Accuracy","Value":metrics['accuracy']},
    {"Metric":"Precision","Value":metrics['precision']},
    {"Metric":"Recall","Value":metrics['recall']},
    {"Metric":"F1","Value":metrics['f1']},
    {"Metric":"ROC-AUC","Value":metrics['roc']},
    {"Metric":"PR-AUC","Value":metrics['pr_auc']},
    {"Metric":"FPR","Value":metrics['fpr']},
])
display(df_summary)

# ---------------------------
# 10) Save artifacts
# ---------------------------
joblib.dump(logic_tree, LOGIC_TREE_PATH)
joblib.dump(scaler, SCALER_PATH)
print("Saved logic tree and scaler.")

# Save a small audit CSV
audit_df = pd.DataFrame({
    "session_id": df.iloc[y_test.index if hasattr(y_test, 'index') else range(len(y_test))].get('session_id', range(len(y_test))),
    "y_true": y_test.values,
    "lstm_proba": y_proba_lstm,
    "tree_score": [logic_tree.score_path(p) for p in token_paths_test],
    "combined_score": combined_scores,
    "y_pred": y_pred_hybrid
})
audit_df.to_csv(MODEL_DIR/"logic_tree_audit_log.csv", index=False)
print("Audit log saved:", MODEL_DIR/"logic_tree_audit_log.csv")

# ---------------------------
# 11) Report (pptx + pdf) quick exporter
# ---------------------------
os.makedirs("report_assets", exist_ok=True)
# save plots to report_assets
plt.figure(); sns.heatmap(cm, annot=True, fmt='d', cmap='Blues'); plt.title("Confusion Matrix"); plt.savefig("report_assets/confusion_matrix.png"); plt.close()
plt.figure(); plt.plot(fpr,tpr); plt.title("ROC Curve"); plt.savefig("report_assets/roc_curve.png"); plt.close()
plt.figure(); plt.plot(rec,prec); plt.title("PR Curve"); plt.savefig("report_assets/pr_curve.png"); plt.close()
plt.figure(figsize=(6,4)); sns.barplot(data=df_summary, x="Metric", y="Value"); plt.ylim(0,1); plt.title("KPI Summary"); plt.savefig("report_assets/kpi.png"); plt.close()

# PPTX
prs = Presentation()
slide = prs.slides.add_slide(prs.slide_layouts[0])
slide.shapes.title.text = "Logic Tree + LSTM Hybrid System Report"
slide.placeholders[1].text = "Incognoir Cybertech | AI + Security Evaluation"

# KPI Table for PPTX
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = "Model KPI Summary"
# Prepare data as a list of lists for PPTX table, ensuring numeric values are rounded
pptx_table_data = [df_summary.columns.tolist()] # Header row
for index, row in df_summary.iterrows():
    metric_name = row['Metric']
    metric_value = row['Value']
    # Explicitly handle potential non-numeric values in the 'Value' column
    if isinstance(metric_value, (int, float, np.number)) and not pd.isna(metric_value):
        value_text = str(round(float(metric_value), 3)) # Round numeric values
    else:
        value_text = str(metric_value) # Keep as string (e.g., 'nan') if not numeric
    pptx_table_data.append([str(metric_name), value_text])

rows, cols = len(pptx_table_data), len(pptx_table_data[0])
table = slide.shapes.add_table(rows, cols, Inches(1), Inches(1.5), Inches(8), Inches(4)).table

# Populate the table cells
for r_idx, row_data in enumerate(pptx_table_data):
    for c_idx, cell_data in enumerate(row_data):
        table.cell(r_idx, c_idx).text = str(cell_data) # Ensure data is string


# Add Charts
for fname in sorted(Path("report_assets").glob("*.png")):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.shapes.add_picture(str(fname), Inches(1), Inches(1.5), Inches(8), Inches(4.5))
prs.save("LogicTree_LSTM_Report.pptx")
print("Saved PPTX:", REPORT_PPTX)

# PDF
doc = SimpleDocTemplate(REPORT_PDF, pagesize=A4)
styles = getSampleStyleSheet()
story = [Paragraph("Logic Tree + LSTM Evaluation Report", styles["Title"])]
# Add KPI Summary table to PDF
# Prepare data as a list of lists, explicitly handling rounding for the 'Value' column
data = [df_summary.columns.tolist()] # Add header row
for index, row in df_summary.iterrows():
    # Explicitly convert Value to float before rounding and then to string
    value = float(row['Value']) if pd.notna(row['Value']) else 0.0
    data.append([row['Metric'], str(round(value, 3))]) # Append [Metric, rounded Value as string]

table = Table(data)
# Add table style (optional)
style = TableStyle([
    ('BACKGROUND', (0,0), (-1,0), colors.grey),
    ('TEXTCOLOR', (0,0), (-1,0), colors.whitesmoke),
    ('ALIGN', (0,0), (-1,-1), 'CENTER'),
    ('FONTNAME', (0,0), (-1,0), 'Helvetica-Bold'),
    ('BOTTOMPADDING', (0,0), (-1,0), 12),
    ('BACKGROUND', (0,1), (-1,-1), colors.beige),
    ('GRID', (0,0), (-1,-1), 1, colors.black)
])
table.setStyle(style)
story.append(table)
story.append(Spacer(1, 12))


for f in os.listdir("report_assets"):
    if f.endswith(".png"):
        story.append(RLImage(f"report_assets/{f}", width=400, height=300))
doc.build(story)

print("\n🎯 Final Deliverables Generated:")
print("• logic_tree_audit_log.csv")
print("• LogicTree_LSTM_Report.pptx")
print("• LogicTree_LSTM_Report.pdf")