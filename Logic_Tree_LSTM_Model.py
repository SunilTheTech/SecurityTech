# ===============================================================
# 🧾 Final Report Generator for Logic Tree + LSTM
# ===============================================================
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Image as RLImage
from reportlab.lib.styles import getSampleStyleSheet
from reportlab.lib.pagesizes import A4
import os
import matplotlib.pyplot as plt # Import matplotlib
import seaborn as sns # Import seaborn
from sklearn.metrics import confusion_matrix, roc_curve, precision_recall_curve, auc # Import metrics for plots

# Ensure directories
os.makedirs("report_assets", exist_ok=True)

# Save charts as images
# Ensure df_summary is available from previous cells
if 'df_summary' in globals():
    # Re-generate plots if they were not saved or figures were closed
    # This is a safeguard; ideally, the previous cells generating plots are run
    # before this cell.

    # Plot Accuracy (for applicable models)
    plt.figure(figsize=(10, 6))
    # Filter out Isolation Forest for Accuracy plot
    # Ensure metrics_df_comparison is available from previous cells
    if 'metrics_df_comparison' in globals():
        accuracy_plot_df = metrics_df_comparison.dropna(subset=['Accuracy']).copy() # Use .copy() to avoid SettingWithCopyWarning
        plt.bar(accuracy_plot_df['Model'], accuracy_plot_df['Accuracy'], color=['blue', 'green', 'red', 'purple'])
        plt.axhline(y=0.95, color='gray', linestyle='--', label='KPI: Accuracy >= 95%') # KPI line
        plt.ylim(0.9, 1.0) # Set y-axis limits for better comparison around 1
        plt.ylabel('Accuracy')
        plt.title('Model Comparison: Accuracy')
        plt.xticks(rotation=45, ha='right')
        plt.legend()
        plt.tight_layout()
        plt.savefig("report_assets/chart_accuracy.png", dpi=300, bbox_inches="tight")
        plt.close() # Close figure to free memory

        # Plot ROC-AUC
        plt.figure(figsize=(10, 6))
        plt.bar(metrics_df_comparison['Model'], metrics_df_comparison['ROC-AUC'], color=['blue', 'green', 'red', 'orange', 'purple'])
        plt.axhline(y=0.95, color='gray', linestyle='--', label='Desired ROC-AUC (example)') # Example desired level
        plt.ylim(0.6, 1.0) # Adjust y-axis limits as needed
        plt.ylabel('ROC-AUC')
        plt.title('Model Comparison: ROC-AUC')
        plt.xticks(rotation=45, ha='right')
        plt.legend()
        plt.tight_layout()
        plt.savefig("report_assets/chart_roc_auc.png", dpi=300, bbox_inches="tight")
        plt.close() # Close figure

        # Plot FPR
        plt.figure(figsize=(10, 6))
        plt.bar(metrics_df_comparison['Model'], metrics_df_comparison['FPR'], color=['blue', 'green', 'red', 'orange', 'purple'])
        plt.axhline(y=0.03, color='gray', linestyle='--', label='KPI: FPR <= 3% (0.03)') # KPI line
        plt.ylim(0, 0.15) # Adjust y-axis limits as needed
        plt.ylabel('False Positive Rate (FPR)')
        plt.title('Model Comparison: False Positive Rate (FPR)')
        plt.xticks(rotation=45, ha='right')
        plt.legend()
        plt.tight_layout()
        plt.savefig("report_assets/chart_fpr.png", dpi=300, bbox_inches="tight")
        plt.close() # Close figure

        # Plot PR-AUC (specifically for Isolation Forest)
        plt.figure(figsize=(5, 6))
        # Filter for Isolation Forest
        pr_auc_plot_df = metrics_df_comparison.dropna(subset=['PR-AUC (Anomaly)']).copy() # Use .copy()
        if not pr_auc_plot_df.empty:
            plt.bar(pr_auc_plot_df['Model'], pr_auc_plot_df['PR-AUC (Anomaly)'], color=['orange'])
            plt.ylabel('Precision-Recall AUC')
            plt.title('Isolation Forest: Precision-Recall AUC')
            plt.xticks(rotation=45, ha='right')
            plt.tight_layout()
            plt.savefig("report_assets/chart_pr_auc_iforest.png", dpi=300, bbox_inches="tight")
            plt.close() # Close figure
        else:
            print("\nSkipping PR-AUC plot generation for report: No data available for Isolation Forest PR-AUC.")

    else:
        print("Error: metrics_df_comparison not found. Cannot generate comparison plots for report.")


    # Re-generate Logic Tree + LSTM specific plots if df_audit is available
    if 'df_audit' in globals():
        # Confusion Matrix (Logic Tree + LSTM)
        y_true = df_audit["true_label"]
        y_pred = (df_audit["combined_risk"] >= 0.65).astype(int)
        cm = confusion_matrix(y_true, y_pred)
        plt.figure(figsize=(5,4))
        sns.heatmap(cm, annot=True, fmt='d', cmap="Blues", cbar=False)
        plt.title("Confusion Matrix – Logic Tree + LSTM")
        plt.xlabel("Predicted")
        plt.ylabel("True")
        plt.savefig("report_assets/chart_confusion_matrix_lt_lstm.png", dpi=300, bbox_inches="tight")
        plt.close() # Close figure

        # ROC Curve (Combined Model)
        y_score = df_audit["combined_risk"]
        fpr, tpr, _ = roc_curve(y_true, y_score)
        roc_auc = auc(fpr, tpr)
        plt.figure(figsize=(6,5))
        plt.plot(fpr, tpr, color='darkorange', lw=2, label=f'ROC AUC = {roc_auc:.3f}')
        plt.plot([0,1], [0,1], color='gray', linestyle='--')
        plt.xlabel("False Positive Rate")
        plt.ylabel("True Positive Rate")
        plt.title("ROC Curve – Combined Model")
        plt.legend()
        plt.savefig("report_assets/chart_roc_curve_lt_lstm.png", dpi=300, bbox_inches="tight")
        plt.close() # Close figure

        # Precision-Recall Curve (Combined Model)
        prec, rec, _ = precision_recall_curve(y_true, y_score)
        plt.figure(figsize=(6,5))
        plt.plot(rec, prec, color='purple', lw=2)
        plt.xlabel("Recall")
        plt.ylabel("Precision")
        plt.title("Precision-Recall Curve – Combined Model")
        plt.savefig("report_assets/chart_pr_curve_lt_lstm.png", dpi=300, bbox_inches="tight")
        plt.close() # Close figure

        # Overall Sandbox Decision Distribution
        sandbox_rate = (df_audit["action"] == "SANDBOX").mean() * 100
        plt.figure(figsize=(5,5))
        plt.pie(
            [sandbox_rate, 100 - sandbox_rate],
            labels=[f"Sandboxed ({sandbox_rate:.1f}%)", "Allowed"],
            colors=["#ff6666","#66b3ff"],
            autopct="%1.1f%%",
            startangle=140
        )
        plt.title("Overall Sandbox Decision Distribution")
        plt.savefig("report_assets/chart_sandbox_distribution.png", dpi=300, bbox_inches="tight")
        plt.close() # Close figure

        # Risk Evolution Plot
        df_sorted = df_audit.sort_values("combined_risk", ascending=False).reset_index(drop=True)
        plt.figure(figsize=(12,6))
        sns.lineplot(x=range(len(df_sorted)), y=df_sorted["combined_risk"], label="Combined Risk")
        plt.axhline(0.65, color='r', linestyle='--', label="Threshold")
        plt.title("Risk Scores Across Sessions – Logic Tree + LSTM System")
        plt.xlabel("Session Index")
        plt.ylabel("Combined Risk Score")
        plt.legend()
        plt.savefig("report_assets/chart_risk_evolution.png", dpi=300, bbox_inches="tight")
        plt.close() # Close figure

        # KPI Summary Bar Plot (Logic Tree + LSTM)
        # Ensure df_summary is available
        if 'df_summary' in globals():
            plt.figure(figsize=(7,5))
            sns.barplot(data=df_summary, x="Metric", y="Value", palette="viridis")
            plt.title("Model Performance KPIs (Logic Tree + LSTM)")
            plt.ylim(0,1)
            plt.savefig("report_assets/chart_kpi_summary_bar.png", dpi=300, bbox_inches="tight")
            plt.close() # Close figure
        else:
            print("Error: df_summary not found. Cannot generate KPI summary bar plot for report.")

    else:
        print("Error: df_audit not found. Cannot generate Logic Tree + LSTM specific plots for report.")

else:
    print("Error: df_summary not found. Cannot generate any plots for the report.")


# Create PowerPoint
prs = Presentation()
title_slide = prs.slides.add_slide(prs.slide_layouts[0])
if title_slide.shapes.title: # Add check for title placeholder
    title_slide.shapes.title.text = "Logic Tree + LSTM Evaluation Report"
if title_slide.placeholders[1]: # Add check for subtitle placeholder
    title_slide.placeholders[1].text = "Incognoir Cybertech – Deeptech AI Security | Auto-generated summary"

# Add summary text
slide = prs.slides.add_slide(prs.slide_layouts[1])
if slide.shapes.title: # Add check for title placeholder
    slide.shapes.title.text = "System Overview"
body = slide.placeholders[1]
body.text = (
    "This report summarizes the performance of the Logic Tree + LSTM hybrid security model.\n"
    "It includes session-wise risk tracking, audit logs, performance KPIs, and sandbox behavior analytics.\n"
    "Generated using real-time evaluation and visualization pipelines in Google Colab."
)

# KPI Chart slide (Using the df_summary for KPI values)
if 'df_summary' in globals() and not df_summary.empty:
    slide = prs.slides.add_slide(prs.slide_layouts[5]) # Use a layout with a title and content placeholder
    if slide.shapes.title: # Add check for title placeholder
        slide.shapes.title.text = "Logic Tree + LSTM KPI Summary"
    # Add a table for KPI summary instead of a chart for simplicity
    x, y, cx, cy = Inches(1), Inches(1.5), Inches(8), Inches(4)
    shape = slide.shapes.add_table(rows=len(df_summary)+1, cols=2, left=x, top=y, width=cx, height=cy)
    table = shape.table

    # Set column widths (optional)
    table.columns[0].width = Inches(4)
    table.columns[1].width = Inches(4)

    # Write table headers
    table.cell(0, 0).text = "Metric"
    table.cell(0, 1).text = "Value"

    # Write data rows
    for i, row in df_summary.iterrows():
        table.cell(i + 1, 0).text = str(row['Metric'])
        table.cell(i + 1, 1).text = f"{row['Value']:.3f}" # Format value


else:
    print("Skipping KPI Summary Table slide: df_summary not found or empty.")


# Add charts as slides
# Ensure report_assets directory exists and contains images
if os.path.exists("report_assets") and os.listdir("report_assets"):
    # Sort files to ensure consistent order in the report
    chart_files = sorted([f for f in os.listdir("report_assets") if f.endswith(".png")])
    for file in chart_files:
        slide = prs.slides.add_slide(prs.slide_layouts[6]) # Use a blank layout
        # Set title based on filename
        if slide.shapes.title: # Add check for title placeholder
            slide.shapes.title.text = file.replace("chart_", "").replace(".png", "").replace("_", " ").title()
        # Add picture, adjusting position and size
        img_path = os.path.join("report_assets", file)
        # Adjust picture size and position based on slide layout and image aspect ratio if needed
        # For now, use fixed size and position
        slide.shapes.add_picture(img_path, Inches(1), Inches(1.5), Inches(8), Inches(4.5))
else:
    print("Skipping chart slides: report_assets directory not found or is empty.")


pptx_path = "LogicTree_LSTM_Report.pptx"
prs.save(pptx_path)
print(f"✅ PowerPoint saved: {pptx_path}")

# PDF Report
pdf_path = "LogicTree_LSTM_Report.pdf"
doc = SimpleDocTemplate(pdf_path, pagesize=A4)
styles = getSampleStyleSheet()
story = [
    Paragraph("<b>Logic Tree + LSTM Evaluation Report</b>", styles["Title"]),
    Spacer(1, 12),
    Paragraph("Incognoir Cybertech – Deeptech AI Security", styles["Heading2"]),
    Spacer(1, 12),
    Paragraph("This PDF summarizes model performance metrics and key visuals for investor/stakeholder reporting.", styles["BodyText"]),
    Spacer(1, 24)
]

# Add KPI Table to PDF
if 'df_summary' in globals() and not df_summary.empty:
    story.append(Paragraph("<b>Model KPI Summary</b>", styles["Heading3"]))
    story.append(Spacer(1, 12))
    # Using reportlab.platypus.Table for a structured table
    from reportlab.platypus import Table, TableStyle
    from reportlab.lib import colors

    data = [["Metric", "Value"]] # Table header
    for i, row in df_summary.iterrows():
        data.append([str(row['Metric']), f"{row['Value']:.3f}"]) # Add data rows

    table = Table(data)

    # Add table style (optional)
    style = TableStyle([
        ('BACKGROUND', (0, 0), (-1, 0), colors.grey),
        ('TEXTCOLOR', (0, 0), (-1, 0), colors.whitesmoke),
        ('ALIGN', (0, 0), (-1, -1), 'CENTER'),
        ('FONTNAME', (0, 0), (-1, 0), 'Helvetica-Bold'),
        ('BOTTOMPADDING', (0, 0), (-1, 0), 12),
        ('BACKGROUND', (0, 1), (-1, -1), colors.beige),
        ('GRID', (0, 0), (-1, -1), 1, colors.black),
    ])
    table.setStyle(style)

    story.append(table)
    story.append(Spacer(1, 24))
else:
     print("Skipping KPI Summary Table in PDF: df_summary not found or empty.")


# Add figures to PDF
if os.path.exists("report_assets") and os.listdir("report_assets"):
    # Sort files to ensure consistent order in the report
    chart_files = sorted([f for f in os.listdir("report_assets") if f.endswith(".png")])
    for file in chart_files:
        story.append(Spacer(1, 12))
        img_path = os.path.join("report_assets", file)
        # Adjust image width and height for PDF
        img = RLImage(img_path, width=400, height=300) # Adjust size as needed
        story.append(img)
else:
    print("Skipping chart images in PDF: report_assets directory not found or is empty.")


doc.build(story)
print(f"✅ PDF saved: {pdf_path}")

print("\n🎯 Final deliverables generated:")
print(f"• {pptx_path}")
print(f"• {pdf_path}")
print(f"• logic_tree_audit_log.csv (raw audit data, if generated)")
print("All assets are ready for download and presentation.")